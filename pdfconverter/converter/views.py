from django.shortcuts import render, redirect
from django.contrib import messages
from django.contrib.auth import authenticate, login
from django.http import HttpResponse, JsonResponse
from django.views.decorators.csrf import csrf_exempt
from reportlab.pdfgen import canvas
from io import BytesIO
from PyPDF2 import PdfReader
from PIL import Image
import pandas as pd
from docx import Document
import os
from django.conf import settings
from django.http import FileResponse
from django.core.files.storage import default_storage
from PyPDF2 import PdfReader, PdfWriter
from django.shortcuts import render
from django.core.files.base import ContentFile
from reportlab.lib.pagesizes import letter
from zipfile import ZipFile
from PIL import Image
from pdf2image import convert_from_path
import ebooklib
from ebooklib import epub
from pptx import Presentation
from pptx.util import Inches,Pt
import subprocess
from ebooklib import ITEM_DOCUMENT
from pdf2image import convert_from_bytes
import fitz
from bs4 import BeautifulSoup
import tabula
import pdfplumber
from .models import CustomUser
import ssl



def generate_pdf_response(buffer, original_name):
    """Helper function to generate a PDF response with the same name as the uploaded file."""
    base_name = os.path.splitext(original_name)[0]  # Get the base name without extension
    buffer.seek(0)
    response = HttpResponse(buffer, content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="{base_name}.pdf"'
    return response

def generate_ppt_response(buffer, original_name):
    """Helper function to generate a PPTX response with the same name as the uploaded file."""
    base_name = os.path.splitext(original_name)[0]  # Get the base name without extension
    buffer.seek(0)
    response = HttpResponse(buffer, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = f'attachment; filename="{base_name}.pptx"'
    return response


def convert_pdf_to_pptx(uploaded_file):
    """
    Converts a PDF file into a PowerPoint presentation.

    Args:
        uploaded_file: A file object for the uploaded PDF.

    Returns:
        BytesIO buffer containing the PowerPoint file.
    """
    pdf_document = fitz.Document(stream=uploaded_file)  # Open PDF as a stream
    presentation = Presentation()

    # Convert each PDF page to a PowerPoint slide
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        text = page.get_text("text")  # Extract text from PDF page

        # Add a new slide with a blank layout
        slide_layout = presentation.slide_layouts[5]  # Blank slide layout
        slide = presentation.slides.add_slide(slide_layout)

        # Add a textbox to the slide
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        # Add text content line by line
        for line in text.split("\n"):
            if line.strip():  # Ignore empty lines
                paragraph = text_frame.add_paragraph()
                paragraph.text = line.strip()
                paragraph.font.size = Pt(14)

    # Save the presentation to a BytesIO buffer
    buffer = BytesIO()
    presentation.save(buffer)
    buffer.seek(0)  # Move to the start of the buffer
    return buffer


def generate_ppt_response(buffer, original_filename):
    """
    Generates an HTTP response for downloading the PowerPoint file.

    Args:
        buffer: BytesIO object containing the PowerPoint file.
        original_filename: The original uploaded filename.

    Returns:
        HttpResponse object for downloading the file.
    """
    pptx_filename = original_filename.replace(".pdf", ".pptx")
    response = HttpResponse(buffer, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    response["Content-Disposition"] = f'attachment; filename="{pptx_filename}"'
    return response



def convert_pdf_to_png(uploaded_file):
    """
    Converts each page of a PDF file into PNG images.

    Args:
        uploaded_file: A file object for the uploaded PDF.

    Returns:
        A BytesIO object containing a ZIP file with PNG images.
    """
    pdf_document = fitz.Document(stream=uploaded_file)  # Open the PDF file
    zip_buffer = BytesIO()

    # Create a ZIP archive to store PNG images
    with ZipFile(zip_buffer, "w") as zip_archive:
        for page_num in range(len(pdf_document)):
            # Load the page
            page = pdf_document.load_page(page_num)

            # Render the page as a Pixmap (image representation)
            pix = page.get_pixmap(dpi=150)  # Adjust DPI for image quality
            image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            # Save the image to a BytesIO buffer
            image_buffer = BytesIO()
            image.save(image_buffer, format="PNG")
            image_buffer.seek(0)

            # Add the image to the ZIP archive
            zip_archive.writestr(f"page_{page_num + 1}.png", image_buffer.read())

    zip_buffer.seek(0)  # Move to the start of the ZIP buffer
    return zip_buffer

def convert_pdf_to_excel(uploaded_file):
    """
    Converts a PDF file into an Excel file.

    Args:
        uploaded_file: A file object for the uploaded PDF.

    Returns:
        BytesIO buffer containing the Excel file.
    """
    buffer = BytesIO()

    try:
        # Open the PDF file using pdfplumber
        with pdfplumber.open(uploaded_file) as pdf:
            all_data = []
            for page in pdf.pages:
                # Extract tables from each page
                tables = page.extract_tables()
                for table in tables:
                    # Append the table data
                    all_data.append(table)

        # Flatten and write data to Excel
        with pd.ExcelWriter(buffer, engine='xlsxwriter') as writer:
            for i, table in enumerate(all_data):
                if table:
                    # Convert the table data into a DataFrame
                    df = pd.DataFrame(table)
                    sheet_name = f"Sheet_{i+1}"
                    df.to_excel(writer, index=False, header=False, sheet_name=sheet_name)

        buffer.seek(0)  # Move to the start of the buffer
        return buffer
    except Exception as e:
        raise Exception(f"Failed to convert PDF to Excel: {str(e)}")


def home(request):
    return render(request, 'index.html')

def signin(request):
    if request.method == 'POST':
        email = request.POST['email']
        password = request.POST['password']

        # Check if a user with this email already exists
        if CustomUser.objects.filter(email=email).exists():
            # Authenticate the user
            user = authenticate(request, username=email, password=password)
            if user is not None:
                login(request, user)
                messages.success(request, 'Login successful!')
                return redirect('home')  # Redirect to the home page
            else:
                messages.error(request, 'Invalid email or password. Please try again.')
        else:
            # If the user doesn't exist, create a new one
            obj = CustomUser()
            obj.email = email
            obj.password = password
            obj.save()
            messages.success(request, 'Account created successfully! Please log in.')

        return render(request, 'index.html')  # Show the sign-in form with an error message if needed

    return render(request, 'signin.html')

def forgot_password(request):
    if request.method == 'POST':
        # Handle the forgot password logic here
        email = request.POST.get('email')
        # Add logic to send a reset email
    return render(request, 'forgot_password.html')

def contact_us(request):
    return render(request, 'Contact_Us.html')

def add_page_number(request):
    if request.method == 'POST' and 'files' in request.FILES:
        try:
            uploaded_file = request.FILES['files']

            # Validate file type
            if not uploaded_file.name.endswith('.pdf'):
                return JsonResponse({'error': 'Invalid file type. Please upload a PDF file.'}, status=400)

            # Save the uploaded file temporarily
            temp_dir = os.path.join(settings.MEDIA_ROOT, 'temp')
            os.makedirs(temp_dir, exist_ok=True)  # Ensure the directory exists
            temp_file_path = os.path.join(temp_dir, uploaded_file.name)

            with open(temp_file_path, 'wb') as temp_file:
                for chunk in uploaded_file.chunks():
                    temp_file.write(chunk)

            # Read the uploaded PDF
            pdf_reader = PdfReader(temp_file_path)
            pdf_writer = PdfWriter()

            # Create a buffer to hold the modified PDF
            for page_number, page in enumerate(pdf_reader.pages, start=1):
                # Create a new PDF with page number
                packet = BytesIO()
                can = canvas.Canvas(packet, pagesize=letter)
                can.drawString(500, 10, f"Page {page_number}")  # Adjust coordinates (500, 10) for positioning
                can.save()

                # Merge the page number PDF with the original page
                packet.seek(0)
                new_pdf = PdfReader(packet)
                page.merge_page(new_pdf.pages[0])

                pdf_writer.add_page(page)

            # Write the modified PDF to a BytesIO buffer
            buffer = BytesIO()
            pdf_writer.write(buffer)
            buffer.seek(0)

            # Prepare the response for downloading the modified file
            base_name = os.path.splitext(uploaded_file.name)[0]  # File name without extension
            response = HttpResponse(buffer, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{base_name}_with_page_numbers.pdf"'

            # Clean up temporary files
            os.remove(temp_file_path)

            return response

        except Exception as e:
            # Handle any errors
            return JsonResponse({'error': str(e)}, status=400)

    # Render the upload form if not a POST request or no file uploaded
    return render(request, 'Add_page_number.html')

def merge(request):
    if request.method == 'POST' and 'files' in request.FILES:
        try:
            uploaded_files = request.FILES.getlist('files')  # Get all uploaded files

            # Check if the number of files exceeds the limit
            if len(uploaded_files) > 5:
                return JsonResponse({'error': 'You have achieved the maximum amount of files. Please upload up to 5 files only.'}, status=400)

            if not uploaded_files:
                return JsonResponse({'error': 'No files uploaded.'}, status=400)

            pdf_writer = PdfWriter()

            # Validate and process each uploaded PDF
            for uploaded_file in uploaded_files:
                if not uploaded_file.name.endswith('.pdf'):
                    return JsonResponse({'error': f"Invalid file type: {uploaded_file.name}. Please upload only PDF files."}, status=400)
                
                # Read each PDF and add its pages to the writer
                pdf_reader = PdfReader(uploaded_file)
                for page in pdf_reader.pages:
                    pdf_writer.add_page(page)

            # Write the combined PDF to a BytesIO buffer
            buffer = BytesIO()
            pdf_writer.write(buffer)
            buffer.seek(0)

            # Generate a response for the merged PDF
            response = HttpResponse(buffer, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="merged_file.pdf"'

            return response

        except Exception as e:
            return JsonResponse({'error': str(e)}, status=400)

    # Render the merge page for GET requests
    return render(request, 'merge.html')

def remove(request):
    if request.method == 'POST':
        if 'file' not in request.FILES:
            return JsonResponse({'error': 'No file uploaded.'}, status=400)

        if 'pages' not in request.POST or not request.POST['pages']:
            return JsonResponse({'error': 'Please provide pages to remove.'}, status=400)

        uploaded_file = request.FILES['file']
        if not uploaded_file.name.endswith('.pdf'):
            return JsonResponse({'error': 'Invalid file type. Please upload a PDF file.'}, status=400)

        try:
            # Parse page numbers input
            page_numbers = request.POST['pages']
            pages_to_remove = set(int(num.strip()) - 1 for num in page_numbers.split(',') if num.strip().isdigit())

            # Read the uploaded PDF
            pdf_reader = PdfReader(uploaded_file)
            pdf_writer = PdfWriter()

            # Remove specified pages
            for i, page in enumerate(pdf_reader.pages):
                if i not in pages_to_remove:  # Keep the pages not in the removal list
                    pdf_writer.add_page(page)

            # Write the updated PDF to a buffer
            buffer = BytesIO()
            pdf_writer.write(buffer)
            buffer.seek(0)

            # Create a response for the updated PDF
            response = HttpResponse(buffer, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="updated_file.pdf"'
            return response

        except Exception as e:
            return JsonResponse({'error': f"An error occurred: {str(e)}"}, status=500)

    return render(request, 'remove.html')

def rotate(request):
    if request.method == 'POST':
        if 'file' not in request.FILES:
            return JsonResponse({'error': 'No file uploaded.'}, status=400)

        if 'rotation' not in request.POST or not request.POST['rotation']:
            return JsonResponse({'error': 'Please specify the rotation direction and degrees.'}, status=400)

        uploaded_file = request.FILES['file']
        if not uploaded_file.name.endswith('.pdf'):
            return JsonResponse({'error': 'Invalid file type. Please upload a PDF file.'}, status=400)

        try:
            rotation = request.POST['rotation']
            degree = int(request.POST['degrees'])

            # Validate rotation input
            if degree not in [90, 180, 270]:
                return JsonResponse({'error': 'Invalid rotation degrees. Choose 90, 180, or 270.'}, status=400)

            # Read the uploaded PDF
            pdf_reader = PdfReader(uploaded_file)
            pdf_writer = PdfWriter()

            # Apply rotation
            for page in pdf_reader.pages:
                if rotation == 'clockwise':
                    page.rotate(degree)
                elif rotation == 'counterclockwise':
                    page.rotate(-degree)
                pdf_writer.add_page(page)

            # Write the rotated PDF to a buffer
            buffer = BytesIO()
            pdf_writer.write(buffer)
            buffer.seek(0)

            # Return the rotated PDF as a direct download
            response = HttpResponse(buffer, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="rotated_file.pdf"'
            return response

        except Exception as e:
            return JsonResponse({'error': f"An error occurred: {str(e)}"}, status=500)

    return render(request, 'rotate_pdf.html')

def word_to_pdf(request):
    if request.method == 'POST':
        try:
            if 'wordToPdfFile' in request.FILES:
                uploaded_file = request.FILES['wordToPdfFile']
                if uploaded_file.name.endswith('.docx'):
                    doc = Document(uploaded_file)
                    buffer = BytesIO()
                    pdf = canvas.Canvas(buffer)
                    y = 750

                    for para in doc.paragraphs:
                        text = para.text.strip()
                        if text:  # Avoid empty lines
                            pdf.drawString(50, y, text)
                            y -= 20
                            if y < 50:
                                pdf.showPage()
                                y = 750

                    pdf.save()
                    return generate_pdf_response(buffer, uploaded_file.name)
                return JsonResponse({'error': 'Invalid file type. Please upload a Word (.docx) file.'}, status=400)
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)

    return render(request, 'word_to_pdf.html')

def pdf_to_word(request):
    if request.method == 'POST':
        try:
            if 'pdfToWordFile' in request.FILES:
                uploaded_file = request.FILES['pdfToWordFile']
                if uploaded_file.name.endswith('.pdf'):
                    reader = PdfReader(uploaded_file)
                    doc = Document()
                    for page in reader.pages:
                        doc.add_paragraph(page.extract_text())

                    buffer = BytesIO()
                    doc.save(buffer)
                    buffer.seek(0)

                    base_name = os.path.splitext(uploaded_file.name)[0]  # Get the base name without extension
                    response = HttpResponse(buffer, content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                    response['Content-Disposition'] = f'attachment; filename="{base_name}.docx"'
                    return response
                return JsonResponse({'error': 'Invalid file type. Please upload a PDF file.'}, status=400)
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)

    return render(request, 'pdf_to_word.html')

def excel_to_pdf(request):
    if request.method == 'POST':
        try:
            if 'excelToPdfFile' in request.FILES:
                uploaded_file = request.FILES['excelToPdfFile']
                if uploaded_file.name.endswith('.xlsx'):
                    df = pd.read_excel(uploaded_file)
                    buffer = BytesIO()
                    pdf = canvas.Canvas(buffer)
                    y = 750

                    for row in df.itertuples(index=False):
                        pdf.drawString(50, y, str(row))
                        y -= 20
                        if y < 50:
                            pdf.showPage()
                            y = 750

                    pdf.save()
                    return generate_pdf_response(buffer, uploaded_file.name)
                return JsonResponse({'error': 'Invalid file type. Please upload an Excel (.xlsx) file.'}, status=400)
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)

    return render(request, 'excel_to_pdf.html')

def pdf_to_excel(request):
    if request.method == 'POST':
        try:
            if 'pdftoexcelfile' in request.FILES:
                uploaded_file = request.FILES['pdftoexcelfile']
                if uploaded_file.name.endswith('.pdf'):
                    try:
                        excel_buffer = convert_pdf_to_excel(uploaded_file)
                        response = HttpResponse(excel_buffer, content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
                        response['Content-Disposition'] = f'attachment; filename="{uploaded_file.name.replace(".pdf", ".xlsx")}"'
                        return response
                    except Exception as e:
                        return JsonResponse({"error": f"An error occurred: {str(e)}"}, status=500)
                return JsonResponse({'error': 'Invalid file type. Please upload a PDF file.'}, status=400)
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)

    return render(request, 'pdf_to_excel.html')

def ppt_to_pdf(request):
    if request.method == 'POST':
        try:
            if 'powerpointtopdffiles' in request.FILES:
                uploaded_file = request.FILES['powerpointtopdffiles']
                if uploaded_file.name.endswith(('.ppt', '.pptx')):
                    presentation = Presentation(uploaded_file)
                    buffer = BytesIO()
                    pdf = canvas.Canvas(buffer, pagesize=(800, 600))  # Default page size for PowerPoint slides

                    for slide in presentation.slides:
                        y_position = 550
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                pdf.drawString(50, y_position, shape.text.strip())
                                y_position -= 20
                                if y_position < 50:
                                    pdf.showPage()
                                    y_position = 550
                        pdf.showPage()

                    pdf.save()
                    return generate_pdf_response(buffer, uploaded_file.name)
                return JsonResponse({'error': 'Invalid file type. Please upload a PowerPoint (.ppt or .pptx) file.'}, status=400)
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)

    return render(request, 'ppt_to_pdf.html')

def pdf_to_ppt(request):
    if request.method == 'POST':
        try:
            if 'pdftopowerpointfiles' in request.FILES:
                uploaded_file = request.FILES['pdftopowerpointfiles']
                if uploaded_file.name.endswith('.pdf'):
                    try:
                        uploaded_file_data = BytesIO(uploaded_file.read())
                        ppt_buffer = convert_pdf_to_pptx(uploaded_file_data)
                        return generate_ppt_response(ppt_buffer, uploaded_file.name)
                    except Exception as e:
                        return JsonResponse({"error": f"An error occurred: {str(e)}"}, status=500)
                return JsonResponse({'error': 'Invalid file type. Please upload a PDF file.'}, status=400)
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)

    return render(request, 'pdf_to_ppt.html')

def jpg_to_pdf(request):
    if request.method == 'POST':
        try:
            if 'jpgToPdfFile' in request.FILES:
                uploaded_file = request.FILES['jpgToPdfFile']
                if uploaded_file.name.endswith(('.jpg', '.jpeg', '.png')):
                    img = Image.open(uploaded_file)
                    buffer = BytesIO()
                    img.save(buffer, format="PDF")
                    return generate_pdf_response(buffer, uploaded_file.name)
                return JsonResponse({'error': 'Invalid file type. Please upload a JPG or PNG file.'}, status=400)
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)

    return render(request, 'jpg_to_pdf.html')

def pdf_to_png(request):
    if request.method == 'POST':
        if 'pdftopngfiles' in request.FILES:
            uploaded_file = request.FILES['pdftopngfiles']
            if uploaded_file.name.endswith('.pdf'):
                try:
                    uploaded_file_data = BytesIO(uploaded_file.read())
                    zip_buffer = convert_pdf_to_png(uploaded_file_data)
                    response = HttpResponse(zip_buffer, content_type="application/zip")
                    response['Content-Disposition'] = f'attachment; filename="{uploaded_file.name.replace(".pdf", "_images.zip")}"'
                    return response
                except Exception as e:
                    return JsonResponse({"error": f"An error occurred: {str(e)}"}, status=500)
            return JsonResponse({'error': 'Invalid file type. Please upload a PDF file.'}, status=400)

    return render(request, 'pdf_to_png.html')